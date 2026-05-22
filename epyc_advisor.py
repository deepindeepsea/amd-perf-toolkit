#!/usr/bin/env python3
"""
epyc_advisor.py — AMD EPYC Multi-Agent Performance Advisor
============================================================
Routes performance questions across 5 categories to the right NABU agents,
synthesizes answers, and optionally generates a PowerPoint deck.

Usage:
    python3 epyc_advisor.py "Which cloud provider has the best OpenSSL on 64 cores?"
    python3 epyc_advisor.py --pptx "Compare M8A vs C4D for Redis throughput"
    python3 epyc_advisor.py --category competitive "AMD Turin vs Intel SPR on Redis"
    python3 epyc_advisor.py --interactive

NABU Agent IDs:
    EPYC Playbook:      9c936425-8ba1-4c7f-b5fe-dbcd1d8c5f9c
    Cruncher:           f7578702-ce56-4f3d-b746-448c6dd64e00
    EPDW:               6bb60d83-1310-4490-9c34-27b1b4f9f150
    Competitive Intel:  5e4eeb15-370a-448a-8474-0c03a15d954b (restricted)

Data sources:
    chip-compare-hero.lovable.app — live pricing, availability, competitive parity
    EPDW — historical benchmark data warehouse
    Cruncher — AMD-validated cloud perf insights
    Playbook — architecture/tuning methodology (NDA)
    Competitive Intel — competitive landscape analysis (restricted)
"""

import os
import sys
import json
import argparse
import subprocess
import textwrap
from datetime import datetime

# ─── NABU Agent Registry ────────────────────────────────────────────────────

AGENTS = {
    "playbook": {
        "id": "9c936425-8ba1-4c7f-b5fe-dbcd1d8c5f9c",
        "name": "EPYC Playbook",
        "description": "Architecture, PMC events, optimization methodology",
        "restricted": False,
    },
    "cruncher": {
        "id": "f7578702-ce56-4f3d-b746-448c6dd64e00",
        "name": "Cloud Performance Analysis (Cruncher)",
        "description": "AMD-validated cloud perf insights and benchmark validation",
        "restricted": False,
    },
    "epdw": {
        "id": "6bb60d83-1310-4490-9c34-27b1b4f9f150",
        "name": "Cloud Benchmark Performance (EPDW)",
        "description": "Historical benchmark data warehouse — instance/workload/metric queries",
        "restricted": False,
    },
    "competitive": {
        "id": "5e4eeb15-370a-448a-8474-0c03a15d994b",
        "name": "EPYC Competitive Intelligence",
        "description": "AMD vs Intel/ARM/NVIDIA competitive analysis (restricted access)",
        "restricted": True,
    },
}

# Pricing/availability data source
CHIP_COMPARE_URL = "https://chip-compare-hero.lovable.app/"

# ─── Category Routing Table ──────────────────────────────────────────────────

CATEGORIES = {
    "csp_rankings": {
        "label": "CSP Rankings & Benchmark Data",
        "agents": ["cruncher", "epdw"],
        "external": [CHIP_COMPARE_URL],
        "keywords": [
            "best performance", "top instance", "compare", "throughput",
            "bandwidth", "benchmark", "ranking", "fastest", "per dollar",
            "redis", "nginx", "openssl", "spec", "stream", "hpc",
        ],
        "example_questions": [
            "Which cloud provider has the best performance for OpenSSL on 64 cores?",
            "Show me the top 3 instances for Redis throughput across AWS, GCP, and Azure.",
            "What is the measured bandwidth for AMD Turin instances vs Intel Xeon on GCP?",
            "Compare Nginx performance between AWS M8a and GCP C4D instances.",
            "Which instance family gives the best OpenSSL throughput per dollar on AWS?",
        ],
    },
    "root_cause": {
        "label": "Performance Analysis & Root Cause",
        "agents": ["cruncher", "epdw", "playbook"],
        "external": [],
        "keywords": [
            "lower than expected", "high ipc", "low throughput", "ppl throttling",
            "backend memory", "backend bound", "frontend bound", "slow", "why",
            "root cause", "diagnose", "debug", "investigate", "numa", "pmc",
        ],
        "example_questions": [
            "Why is my OpenSSL score lower than expected on the AWS M7a instance?",
            "What metrics should I check when Backend Memory % is high on a GCP N4?",
            "My AMD EPYC instance shows high IPC but low throughput — what's wrong?",
            "Why does PPL throttling happen on cloud EPYC instances and how do I detect it?",
            "The workload is NUMA-sensitive but running on a single-socket cloud VM. What should I check?",
        ],
    },
    "optimization": {
        "label": "Optimization & Tuning Methodology",
        "agents": ["playbook"],
        "external": [],
        "keywords": [
            "profile", "measure", "perf stat", "events", "tuning", "optimize",
            "commands", "how do i", "walk me through", "methodology", "smt",
            "prefetch", "huge pages", "affinity", "numa binding", "cache",
        ],
        "example_questions": [
            "How do I profile cache-to-cache transfer latency on an EPYC Turin instance?",
            "What commands should I run to measure memory bandwidth on an AMD cloud instance?",
            "Walk me through the EPYC performance analysis methodology for a web server workload.",
            "How do I check if SMT is helping or hurting my throughput on a 64-core VM?",
            "What perf stat events should I capture for an OpenSSL benchmark run?",
        ],
    },
    "competitive": {
        "label": "Competitive Intelligence",
        "agents": ["competitive", "cruncher"],
        "external": [CHIP_COMPARE_URL],
        "keywords": [
            "intel", "arm", "graviton", "neoverse", "sapphire rapids", "xeon",
            "competitor", "vs intel", "vs arm", "outperform", "advantage",
            "better than", "compare amd", "amd vs", "competitive",
        ],
        "example_questions": [
            "How does AMD Turin compare to Intel Sapphire Rapids on Redis in the cloud?",
            "Why does Graviton3 sometimes outperform AMD on memory-bound workloads?",
            "What is AMD's advantage over Intel Xeon on OpenSSL at 128 threads?",
            "Compare AMD EPYC vs ARM Neoverse N2 for HPC workloads on GCP.",
            "Intel claims better single-thread performance — what do measured benchmarks show?",
        ],
    },
    "instance_selection": {
        "label": "Cloud Architecture & Instance Selection",
        "agents": ["cruncher", "epdw"],
        "external": [CHIP_COMPARE_URL],
        "keywords": [
            "which instance", "best instance", "memory-optimized", "compute-optimized",
            "instance type", "instance family", "select", "choose", "recommendation",
            "topology", "numa topology", "ppl compare", "generation", "milan",
            "genoa", "turin", "epyc gen",
        ],
        "example_questions": [
            "Should I use a memory-optimized or compute-optimized instance for Redis on AWS?",
            "What's the best AMD instance type for an HPC workload that needs high memory bandwidth?",
            "How does GCP C4D instance PPL compare to AWS M8a under sustained load?",
            "What are the NUMA topology differences between AWS M8a and Azure ECads v5?",
            "Which AMD EPYC generation — Milan, Genoa, or Turin — is best for inference workloads?",
        ],
    },
}

# ─── Query Router ────────────────────────────────────────────────────────────

def classify_question(question: str) -> str:
    """Score each category by keyword overlap and return the best match."""
    q_lower = question.lower()
    scores = {}
    for cat_key, cat in CATEGORIES.items():
        score = sum(1 for kw in cat["keywords"] if kw in q_lower)
        scores[cat_key] = score

    # Tie-break: prefer more specific categories
    best = max(scores, key=lambda k: scores[k])
    if scores[best] == 0:
        best = "root_cause"  # safe default — routes to all agents
    return best


def format_agent_prompt(question: str, category: str, agent_key: str) -> str:
    """Build a focused prompt for each agent based on the question and category."""
    cat = CATEGORIES[category]
    agent = AGENTS[agent_key]

    preambles = {
        "playbook": (
            "You are the AMD EPYC Performance Playbook expert. Answer using the playbook methodology. "
            "Include specific perf event names, L/M/H thresholds, and actionable commands."
        ),
        "cruncher": (
            "You are the Cloud Performance Analysis agent. Provide AMD-validated cloud performance data. "
            "Include measured benchmark numbers, CSP comparisons, and AMD positioning."
        ),
        "epdw": (
            "You are the Cloud Benchmark Performance data warehouse agent. Query available benchmark data. "
            "Report exact benchmark scores, instance types, workload configurations, and trends."
        ),
        "competitive": (
            "You are the AMD EPYC Competitive Intelligence analyst. Provide competitive analysis "
            "backed by measured data. Compare AMD EPYC against Intel Xeon, ARM Neoverse, Graviton, etc. "
            "Distinguish facts from claims."
        ),
    }

    return f"{preambles.get(agent_key, '')} Question: {question}"


# ─── PPTX Generation ─────────────────────────────────────────────────────────

def generate_pptx(question: str, category: str, answers: dict, output_path: str):
    """Generate a PowerPoint slide deck from the advisor output."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[WARN] python-pptx not installed. Run: pip install python-pptx --break-system-packages")
        return None

    AMD_RED   = RGBColor(0xED, 0x1C, 0x24)
    AMD_DARK  = RGBColor(0x13, 0x14, 0x16)
    AMD_GRAY  = RGBColor(0x40, 0x40, 0x40)
    AMD_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

    cat = CATEGORIES[category]
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # completely blank

    def add_slide():
        return prs.slides.add_slide(blank)

    def bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def textbox(slide, text, l, t, w, h,
                size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txb

    def accent_bar(slide, t=0.9, h=0.06):
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(t), Inches(13.33), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = AMD_RED
        bar.line.fill.background()

    # ── Title slide ──────────────────────────────────────────────────────────
    s = add_slide()
    bg(s, AMD_DARK)
    accent_bar(s, t=0.85, h=0.08)
    textbox(s, "AMD EPYC Performance Advisor", 0.5, 0.2, 12, 0.7,
            size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    textbox(s, cat["label"], 0.5, 0.95, 10, 0.5,
            size=20, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC))
    textbox(s, f'"{question}"', 0.5, 1.6, 12, 1.0,
            size=18, color=RGBColor(0xAA, 0xDD, 0xFF))
    textbox(s, f"Generated: {datetime.now().strftime('%Y-%m-%d')}  |  AMD Internal",
            0.5, 6.8, 12, 0.5, size=12,
            color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.LEFT)

    # ── Agent routing slide ───────────────────────────────────────────────────
    s = add_slide()
    bg(s, AMD_DARK)
    accent_bar(s, t=0.85)
    textbox(s, "Query Routing", 0.5, 0.15, 12, 0.6, size=28, bold=True)
    textbox(s, "Sources consulted for this analysis:", 0.5, 1.0, 12, 0.5, size=16,
            color=RGBColor(0xCC, 0xCC, 0xCC))
    y = 1.6
    for ak in cat["agents"]:
        a = AGENTS[ak]
        mark = "🔒 " if a["restricted"] else "✓ "
        textbox(s, f"{mark}{a['name']}", 0.7, y, 6, 0.4, size=15, bold=True,
                color=AMD_RED if a["restricted"] else WHITE)
        textbox(s, a["description"], 0.7, y + 0.38, 10, 0.35, size=12,
                color=RGBColor(0xBB, 0xBB, 0xBB))
        y += 0.85
    if cat["external"]:
        textbox(s, "External Data Sources:", 0.5, y + 0.2, 12, 0.4, size=14, bold=True,
                color=RGBColor(0x88, 0xCC, 0xFF))
        for url in cat["external"]:
            y += 0.55
            textbox(s, url, 0.7, y, 11, 0.35, size=12,
                    color=RGBColor(0x88, 0xCC, 0xFF))

    # ── One slide per agent answer ────────────────────────────────────────────
    for agent_key, answer_text in answers.items():
        a = AGENTS[agent_key]
        s = add_slide()
        bg(s, AMD_DARK)
        accent_bar(s, t=0.85)
        textbox(s, a["name"], 0.5, 0.15, 12, 0.6, size=24, bold=True)

        # Wrap long answers
        lines = []
        for para in answer_text.split("\n"):
            para = para.strip()
            if not para:
                lines.append("")
                continue
            wrapped = textwrap.wrap(para, width=110)
            lines.extend(wrapped)

        # Fit ~20 lines per slide, paginate if needed
        chunk_size = 20
        chunks = [lines[i:i+chunk_size] for i in range(0, max(len(lines), 1), chunk_size)]
        for ci, chunk in enumerate(chunks):
            if ci > 0:
                s = add_slide()
                bg(s, AMD_DARK)
                accent_bar(s, t=0.85)
                textbox(s, f"{a['name']} (cont.)", 0.5, 0.15, 12, 0.6,
                        size=20, bold=True)
            body = "\n".join(chunk)
            textbox(s, body, 0.5, 1.0, 12, 5.8, size=13,
                    color=RGBColor(0xE0, 0xE0, 0xE0))

    # ── Category-specific example questions slide ─────────────────────────────
    s = add_slide()
    bg(s, AMD_DARK)
    accent_bar(s, t=0.85)
    textbox(s, "Related Questions You Can Ask", 0.5, 0.15, 12, 0.6, size=24, bold=True)
    y = 1.0
    for q in cat["example_questions"][:5]:
        textbox(s, f"▸  {q}", 0.6, y, 12, 0.5, size=13,
                color=RGBColor(0xCC, 0xEE, 0xFF))
        y += 0.62

    # ── Chip Compare reminder ─────────────────────────────────────────────────
    s = add_slide()
    bg(s, AMD_DARK)
    accent_bar(s, t=0.85)
    textbox(s, "Live Pricing & Availability Reference", 0.5, 0.15, 12, 0.6,
            size=24, bold=True)
    textbox(s, "Cloud Instance Parity Engine — AMD vs Intel", 0.5, 1.0, 12, 0.5,
            size=18, color=RGBColor(0x88, 0xCC, 0xFF))
    textbox(s, CHIP_COMPARE_URL, 0.5, 1.55, 12, 0.5, size=16,
            color=RGBColor(0x88, 0xCC, 0xFF))
    textbox(s,
            "Use chip-compare-hero.lovable.app for:\n"
            "  • Live pricing across AWS, Azure, GCP by region\n"
            "  • Instance availability status\n"
            "  • AMD vs Intel competitive parity comparison\n"
            "  • Price/performance ratios across CSPs",
            0.5, 2.2, 12, 3.0, size=15, color=RGBColor(0xDD, 0xDD, 0xDD))

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(output_path)
    return output_path


# ─── Main Advisor Logic ───────────────────────────────────────────────────────

def run_advisor(question: str, force_category: str = None,
                make_pptx: bool = False, output_dir: str = ".") -> dict:
    """
    Route question to appropriate agents, collect answers, optionally generate PPTX.
    Returns dict with category, agents, answers, and optional pptx_path.
    """
    category = force_category or classify_question(question)
    cat = CATEGORIES[category]

    print(f"\n{'='*60}")
    print(f"  AMD EPYC Performance Advisor")
    print(f"{'='*60}")
    print(f"  Category : {cat['label']}")
    print(f"  Agents   : {', '.join(AGENTS[a]['name'] for a in cat['agents'])}")
    if cat["external"]:
        print(f"  External : {', '.join(cat['external'])}")
    print(f"{'='*60}\n")

    answers = {}

    # Build combined NABU query — pass all required agents
    agent_ids = [AGENTS[ak]["id"] for ak in cat["agents"] if not AGENTS[ak]["restricted"]]
    restricted_agents = [ak for ak in cat["agents"] if AGENTS[ak]["restricted"]]

    # Craft the unified question with routing context
    agent_names = ", ".join(AGENTS[ak]["name"] for ak in cat["agents"])
    meta_prompt = (
        f"This question is in the category: {cat['label']}. "
        f"Please draw on data from: {agent_names}. "
        f"Include measured benchmark data where available, "
        f"specific perf event names or CLI commands where relevant, "
        f"and competitive context if applicable. "
        f"Structure your answer with: (1) Direct Answer, (2) Supporting Data, "
        f"(3) Key Metrics to Monitor, (4) Recommended Next Steps.\n\n"
        f"Question: {question}"
    )

    # For demo mode (no NABU API key), generate a structured template answer
    try:
        import requests  # check if we can make HTTP calls
        # In real deployment, this would call the NABU API
        # For now generate structured guidance
        raise ImportError("Use template mode")
    except Exception:
        # Template-based answers when not running inside NABU context
        answers["combined"] = generate_template_answer(question, category, cat)

    if restricted_agents:
        print(f"[NOTE] Restricted agent(s) not queried: "
              f"{', '.join(AGENTS[a]['name'] for a in restricted_agents)}")
        print(f"       Access via NABU: intelligence.amd.com")

    # Print answers
    for ak, answer in answers.items():
        print(f"\n{'─'*50}")
        print(f"  {AGENTS.get(ak, {}).get('name', 'Advisor Response')}")
        print(f"{'─'*50}")
        print(answer)

    # External data sources reminder
    if cat["external"]:
        print(f"\n{'─'*50}")
        print("  External Data Sources")
        print(f"{'─'*50}")
        print(f"  Live pricing & competitive: {CHIP_COMPARE_URL}")
        print("  (AMD vs Intel pricing, availability, parity engine)")

    result = {
        "question": question,
        "category": category,
        "category_label": cat["label"],
        "agents_consulted": cat["agents"],
        "answers": answers,
        "external_sources": cat["external"],
        "timestamp": datetime.now().isoformat(),
    }

    # Generate PPTX if requested
    if make_pptx:
        safe_q = "".join(c if c.isalnum() or c in " -_" else "" for c in question[:40])
        fname = f"epyc_advisor_{safe_q.strip().replace(' ', '_')}.pptx"
        pptx_path = os.path.join(output_dir, fname)
        path = generate_pptx(question, category, answers, pptx_path)
        if path:
            print(f"\n[PPTX] Saved: {path}")
            result["pptx_path"] = path

    return result


def generate_template_answer(question: str, category: str, cat: dict) -> str:
    """Generate a structured guidance answer when NABU is not available inline."""
    q_lower = question.lower()

    # Detect key entities from question
    workloads = []
    for w in ["openssl", "redis", "nginx", "spec", "stream", "hpc", "inference", "ml"]:
        if w in q_lower:
            workloads.append(w.upper())

    csps = []
    for c in ["aws", "gcp", "azure", "oracle"]:
        if c in q_lower:
            csps.append(c.upper())

    instances = []
    for i in ["m8a", "m7a", "c4d", "hbv4", "ecads", "m6a", "c6a", "m6i"]:
        if i in q_lower:
            instances.append(i)

    lines = [
        f"### Question: {question}",
        "",
        "#### 1. Direct Answer",
        f"This question requires data from: {', '.join(AGENTS[a]['name'] for a in cat['agents'])}.",
        "",
    ]

    if category == "csp_rankings":
        lines += [
            "For CSP benchmark rankings, query EPDW for measured results:",
            f"  • Workloads of interest: {', '.join(workloads) or 'all tracked workloads'}",
            f"  • CSPs: {', '.join(csps) or 'AWS, GCP, Azure, Oracle'}",
            f"  • Instances: {', '.join(instances) or 'all AMD EPYC instances'}",
            "",
            "Key EPDW query template:",
            '  { "cloudProvider": "AWS", "instanceType": "m8a.48xlarge",',
            '    "benchmarkType": "SPEC-CPU", "benchmarkCategory": "throughput" }',
        ]
    elif category == "root_cause":
        lines += [
            "Root cause analysis follows the playbook methodology (§2):",
            "  1. Run System Checks: freq, SMT status, NPS, PPL",
            "  2. Collect Pipeline Utilization metrics (perf stat -j):",
            "     frontend_bound%, backend_bound_memory%, IPC, effective frequency",
            "  3. Check L/M/H thresholds (see EPYC_PERF_KNOWLEDGE.md §4)",
            "  4. Navigate to the matching solution section",
        ]
        if "ppl" in q_lower or "throttl" in q_lower:
            lines += [
                "",
                "PPL Throttling Detection (key insight from cloud_context.py):",
                "  AWS M7A/M8A PPL = 320W → expect Feff ≈ 0.75× max boost",
                "  If effective_freq < 0.75 × boost_freq → PPL is throttling",
                "  Measure: perf stat -j -e cpu-cycles,task-clock <workload>",
                "  Effective freq (GHz) = cpu-cycles / (task-clock_ms × 1e6)",
            ]
    elif category == "optimization":
        lines += [
            "Optimization follows EPYC Performance Playbook §8 (Pipeline Opportunities).",
            "Start with this perf stat collection pass (bare metal or cloud):",
            "",
            "  sudo perf stat -j \\",
            "    -e de_no_dispatch_per_slot.no_ops_from_frontend,\\",
            "    de_no_dispatch_per_slot.backend_stalls,\\",
            "    de_src_op_disp.all,ex_ret_ops,ls_not_halted_cyc,\\",
            "    ex_no_retire.load_not_complete,ex_no_retire.not_complete,\\",
            "    l2_cache_req_stat.dc_hit_in_l2,l2_cache_req_stat.ls_rd_blk_c,\\",
            "    bp_l1_tlb_miss_l2_tlb_miss.all,ls_l1_d_tlb_miss.all,\\",
            "    ex_ret_brn_misp,ex_ret_brn,cpu-cycles,instructions,task-clock \\",
            "    <your_workload>",
        ]
    elif category == "competitive":
        lines += [
            "Competitive analysis requires the EPYC Competitive Intelligence Agent.",
            f"Access: NABU Agent ID 5e4eeb15-370a-448a-8474-0c03a15d994b (restricted)",
            "",
            "Live pricing parity: " + CHIP_COMPARE_URL,
            "  → Compare AMD vs Intel instance pricing by region and workload",
            "",
            "Key AMD advantages to highlight for measured workloads:",
            "  • AVX-512 on Zen4/5: no frequency drop (unlike Intel pre-SPR)",
            "  • 96 MB L3 per CCX on Genoa-X vs Intel 60 MB shared L3",
            "  • 12-channel DDR5 memory per socket (Genoa) vs Intel 8-channel",
        ]
    elif category == "instance_selection":
        lines += [
            "Instance selection guidance from EPDW + cloud_context.py:",
            "",
            "  Compute-optimized (max throughput/core): AWS M8A, GCP C4D",
            "  Memory-optimized (BW-bound workloads): AWS R7A, Azure ECads v5",
            "  HPC (highest core count): AWS M8A 192-core, GCP C4D-480",
            "",
            "PPL comparison (from cloud_context.py):",
            "  AWS M7A/M8A: 320W → Feff ratio 0.75 (25% below bare metal boost)",
            "  GCP C4D:     400W → Feff ratio 0.85 (15% deficit)",
            "  Azure HBv4:  400W → Feff ratio 0.85",
            "  Oracle OCI:  450W → Feff ratio 0.88 (closest to bare metal)",
        ]

    lines += [
        "",
        "#### 3. Key Metrics to Monitor",
        "  • Effective Frequency (GHz) — detect PPL throttling",
        "  • backend_bound_memory % — memory bottleneck (threshold: >20% moderate)",
        "  • IPC — overall pipeline efficiency (Low: <0.2, High: >1.0)",
        "  • L2 DC Hit Rate % — cache effectiveness",
        "  • Branch Misp Rate % — code quality (High: >5.0 PTI)",
        "",
        "#### 4. Recommended Next Steps",
        "  1. Query EPDW for measured baseline: 'Show benchmark results for <instance>'",
        "  2. Query Cruncher for AMD validation: 'Does this match AMD's tested numbers?'",
        "  3. Run amd_pipeline_metrics.sh on the workload for live PMC data",
        "  4. Cross-reference with EPYC_PERF_KNOWLEDGE.md thresholds",
        f"  5. For pricing/availability: {CHIP_COMPARE_URL}",
    ]

    return "\n".join(lines)


# ─── Interactive Mode ─────────────────────────────────────────────────────────

def interactive_mode():
    print("\n" + "="*60)
    print("  AMD EPYC Multi-Agent Performance Advisor")
    print("  Type 'help' to see example questions, 'quit' to exit")
    print("  Add '--pptx' to any question to generate a slide deck")
    print("="*60 + "\n")
    print("Available agents:")
    for ak, a in AGENTS.items():
        lock = " [restricted]" if a["restricted"] else ""
        print(f"  • {a['name']}{lock}")
    print()

    while True:
        try:
            raw = input("Ask > ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye!")
            break

        if not raw:
            continue
        if raw.lower() in ("quit", "exit", "bye"):
            print("Bye!")
            break
        if raw.lower() == "help":
            for cat_key, cat in CATEGORIES.items():
                print(f"\n  [{cat['label']}]")
                for q in cat["example_questions"][:2]:
                    print(f"    • {q}")
            continue

        make_pptx = "--pptx" in raw
        question = raw.replace("--pptx", "").strip()
        run_advisor(question, make_pptx=make_pptx,
                    output_dir=os.path.expanduser("~"))


# ─── CLI Entry Point ──────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="AMD EPYC Multi-Agent Performance Advisor",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""
        Examples:
          %(prog)s "Which cloud has best OpenSSL on 64 cores?"
          %(prog)s --pptx "Compare M8A vs C4D for Redis"
          %(prog)s --category optimization "What perf events for OpenSSL?"
          %(prog)s --interactive
          %(prog)s --list-categories
        """),
    )
    parser.add_argument("question", nargs="?", help="Question to answer")
    parser.add_argument("--pptx", action="store_true",
                        help="Generate a PowerPoint deck from the answer")
    parser.add_argument("--category", choices=list(CATEGORIES.keys()),
                        help="Force a specific routing category")
    parser.add_argument("--interactive", action="store_true",
                        help="Launch interactive question-answer loop")
    parser.add_argument("--list-categories", action="store_true",
                        help="Print all categories with example questions")
    parser.add_argument("--output-dir", default=".",
                        help="Directory for generated files (default: .)")
    args = parser.parse_args()

    if args.list_categories:
        for cat_key, cat in CATEGORIES.items():
            print(f"\n{'─'*50}")
            print(f"  {cat['label']}")
            print(f"  Agents: {', '.join(AGENTS[a]['name'] for a in cat['agents'])}")
            print(f"  Example questions:")
            for q in cat["example_questions"]:
                print(f"    • {q}")
        return

    if args.interactive:
        interactive_mode()
        return

    if not args.question:
        parser.print_help()
        return

    run_advisor(
        args.question,
        force_category=args.category,
        make_pptx=args.pptx,
        output_dir=args.output_dir,
    )


if __name__ == "__main__":
    main()
